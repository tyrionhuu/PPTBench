# pptbench/extractors/notes_extractor.py

from pptx.slide import NotesSlide
from pptbench.utils import unit_conversion

class NotesExtractor:
    def __init__(self, notes_slide: NotesSlide, measurement_unit: str = "pt"):
        self._notes_slide = notes_slide
        self._measurement_unit = measurement_unit

    def extract_text(self) -> str:
        if self._notes_slide and self._notes_slide.notes_text_frame:
            return self._notes_slide.notes_text_frame.text.strip()
        return ""

    def extract_font_info(self) -> list:
        """
        Extracts font information from all text runs within the notes slide.

        Returns:
            A list of dictionaries containing paragraph index, run index, text, font name, and font size.
        """
        if not (self._notes_slide and self._notes_slide.notes_text_frame):
            return []

        font_details = []
        for p_idx, paragraph in enumerate(self._notes_slide.notes_text_frame.paragraphs):
            for r_idx, run in enumerate(paragraph.runs):
                font = run.font
                font_name = font.name if font.name else "Default"
                # Extract font size in points without conversion
                font_size = font.size.pt if font.size else 12  # Default size 12pt
                font_details.append(
                    {
                        "paragraph_index": p_idx,
                        "run_index": r_idx,
                        "text": run.text,
                        "font_name": font_name,
                        "font_size": font_size,  # Size in points
                    }
                )
        return font_details

    def extract_notes(self) -> dict:
        notes_data = {}
        text = self.extract_text()
        if text:
            notes_data["text"] = text
            font_info = self.extract_font_info()
            if font_info:
                notes_data["font_details"] = font_info
        return notes_data
