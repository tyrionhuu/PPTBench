# pptbench/extractors/shape_extractors.py

from pptbench.utils import unit_conversion
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture


class BaseShapeExtractor:
    def __init__(self, shape: BaseShape, measurement_unit: str = "pt"):
        self._shape = shape
        self._measurement_unit = measurement_unit

    def extract_shape_type(self) -> str:
        shape_type = self._shape.shape_type
        if isinstance(shape_type, MSO_SHAPE_TYPE):
            return shape_type.name
        return str(shape_type)

    def extract_height(self) -> int | float:
        return unit_conversion(self._shape.height, self._measurement_unit)

    def extract_width(self) -> int | float:
        return unit_conversion(self._shape.width, self._measurement_unit)

    def extract_left(self) -> int | float:
        return unit_conversion(self._shape.left, self._measurement_unit)

    def extract_top(self) -> int | float:
        return unit_conversion(self._shape.top, self._measurement_unit)

    def set_measurement_unit(self, unit: str) -> None:
        self._measurement_unit = unit

    def extract_shape(self) -> dict:
        return {
            "name": self._shape.name,
            "shape_id": self._shape.shape_id,
            "shape_type": self.extract_shape_type(),
            "measurement_unit": self._measurement_unit,
            "height": self.extract_height(),
            "width": self.extract_width(),
            "left": self.extract_left(),
            "top": self.extract_top(),
        }


class BaseAutoShapeExtractor(BaseShapeExtractor):
    def __init__(self, shape: AutoShape, measurement_unit="pt"):
        super().__init__(shape, measurement_unit)

    def extract_text(self) -> str:
        if self._shape.has_text_frame:
            return self._shape.text  # type: ignore[attr-defined]
        raise AttributeError("Shape does not have a text frame")

    def extract_font_info(self) -> list:
        """
        Extracts font information from all text runs within the shape.

        Returns:
            A list of dictionaries containing paragraph index, run index, text, font name, and font size.
        """
        if not self._shape.has_text_frame:
            raise AttributeError("Shape does not have a text frame")

        font_details = []
        for p_idx, paragraph in enumerate(self._shape.text_frame.paragraphs):
            for r_idx, run in enumerate(paragraph.runs):
                font = run.font
                font_name = font.name if font.name else "Default"
                # Extract font size in points without conversion
                font_size = font.size.pt if font.size else 12  # Default size 12pt
                font_details.append(
                    {
                        "paragraph_index": p_idx,
                        "run_index": r_idx,
                        "text": run.text,
                        "font_name": font_name,
                        "font_size": font_size,  # Size in points
                    }
                )
        return font_details

    def extract_shape(self) -> dict:
        shape_data = super().extract_shape()
        if self._shape.has_text_frame:
            shape_data["text"] = self.extract_text()
            shape_data["font_details"] = self.extract_font_info()
        return shape_data


class PlaceholderExtractor(BaseAutoShapeExtractor):
    def __init__(self, shape: AutoShape, measurement_unit: str = "pt"):
        super().__init__(shape, measurement_unit)

    def extract_placeholder_format(self) -> str:
        placeholder_format = self._shape.placeholder_format
        if hasattr(placeholder_format, "type"):
            placeholder_type = placeholder_format.type
            if isinstance(placeholder_type, PP_PLACEHOLDER_TYPE):
                return placeholder_type.name
        raise AttributeError("Unknown placeholder format")

    def extract_shape(self) -> dict:
        shape_data = super().extract_shape()
        shape_data["placeholder_type"] = self.extract_placeholder_format()
        return shape_data


class FreeformExtractor(BaseAutoShapeExtractor):
    def __init__(self, shape: AutoShape, measurement_unit: str = "pt"):
        super().__init__(shape, measurement_unit)


class ConnectorExtractor(BaseShapeExtractor):
    def __init__(self, shape: Connector, measurement_unit: str = "pt"):
        super().__init__(shape, measurement_unit)

    def extract_begin_x(self) -> int | float:
        return unit_conversion(self._shape.begin_x, self._measurement_unit)  # type: ignore[attr-defined]

    def extract_begin_y(self) -> int | float:
        return unit_conversion(self._shape.begin_y, self._measurement_unit)  # type: ignore[attr-defined]

    def extract_end_x(self) -> int | float:
        return unit_conversion(self._shape.end_x, self._measurement_unit)  # type: ignore[attr-defined]

    def extract_end_y(self) -> int | float:
        return unit_conversion(self._shape.end_y, self._measurement_unit)  # type: ignore[attr-defined]

    def extract_shape(self) -> dict:
        shape_data = super().extract_shape()
        shape_data["begin_x"] = self.extract_begin_x()
        shape_data["begin_y"] = self.extract_begin_y()
        shape_data["end_x"] = self.extract_end_x()
        shape_data["end_y"] = self.extract_end_y()
        return shape_data


class PictureExtractor(BaseShapeExtractor):
    def __init__(self, shape: Picture, measurement_unit: str = "pt"):
        super().__init__(shape, measurement_unit)

    def extract_auto_shape_type(self) -> str | None:
        auto_shape_type = self._shape.auto_shape_type  # type: ignore[attr-defined]
        if isinstance(auto_shape_type, MSO_AUTO_SHAPE_TYPE):
            return auto_shape_type.name
        return None

    def extract_filename(self) -> str | None:
        return self._shape.image.filename  # type: ignore[attr-defined]

    # def _extract_blob_str(self) -> str:
    #     blob = self._shape.image.blob  # type: ignore[attr-defined]
    #     return base64.b64encode(blob)

    def extract_shape(self) -> dict:
        shape_data = super().extract_shape()
        if self.extract_auto_shape_type() is not None:
            shape_data["auto_shape_type"] = self.extract_auto_shape_type()
        # shape_data["blob_str"] = self._extract_blob_str()
        return shape_data


class MovieExtractor(BaseShapeExtractor):
    def __init__(self, shape: Movie, measurement_unit: str = "pt"):
        super().__init__(shape, measurement_unit)


class GraphicFrameExtractor(BaseShapeExtractor):
    def __init__(self, shape: GraphicFrame, measurement_unit: str = "pt"):
        super().__init__(shape, measurement_unit)

    def extract_shape(self) -> dict:
        shape_data = super().extract_shape()
        shape_data["has_chart"] = self._shape.has_chart
        shape_data["has_table"] = self._shape.has_table
        return shape_data


class GroupShapeExtractor(BaseShapeExtractor):
    def __init__(self, shape: GroupShape, measurement_unit: str = "pt"):
        super().__init__(shape, measurement_unit)

    def extract_group_shapes(self) -> list:
        from .factories import (
            shape_extractor_factory,  # Local import to avoid circular import
        )

        group_shape_data = []

        for nested_shape in self._shape.shapes:  # type: ignore[attr-defined]
            extractor = shape_extractor_factory(nested_shape, self._measurement_unit)
            shape_data = extractor.extract_shape()
            group_shape_data.append(shape_data)

        return group_shape_data
