# pptbench/extractors/factories.py

from typing import TypeAlias, Union
import logging

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.autoshape import Shape as AutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Movie, Picture

from .shape_extractors import (
    BaseAutoShapeExtractor,
    BaseShapeExtractor,
    ConnectorExtractor,
    FreeformExtractor,
    GraphicFrameExtractor,
    GroupShapeExtractor,
    MovieExtractor,
    PictureExtractor,
    PlaceholderExtractor,
)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

ShapeExtractor: TypeAlias = Union[
    BaseShapeExtractor,
    BaseAutoShapeExtractor,
    ConnectorExtractor,
    FreeformExtractor,
    GroupShapeExtractor,
    PictureExtractor,
    MovieExtractor,
    PlaceholderExtractor,
]

Shape: TypeAlias = Union[
    BaseShape,
    AutoShape,
    Connector,
    GraphicFrame,
    GroupShape,
    Picture,
    Movie,
]

SHAPE_EXTRACTOR_MAP = {
    # Auto Shape
    MSO_SHAPE_TYPE.AUTO_SHAPE: BaseAutoShapeExtractor,
    MSO_SHAPE_TYPE.TEXT_BOX: BaseAutoShapeExtractor,
    MSO_SHAPE_TYPE.FREEFORM: FreeformExtractor,
    MSO_SHAPE_TYPE.PLACEHOLDER: PlaceholderExtractor,
    # Graphic Frame
    MSO_SHAPE_TYPE.CHART: GraphicFrameExtractor,
    MSO_SHAPE_TYPE.TABLE: GraphicFrameExtractor,
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: GraphicFrameExtractor,
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: GraphicFrameExtractor,
    # Picture
    MSO_SHAPE_TYPE.PICTURE: PictureExtractor,
    MSO_SHAPE_TYPE.MEDIA: MovieExtractor,
    # Connector
    MSO_SHAPE_TYPE.LINE: ConnectorExtractor,
    # Group Shape
    MSO_SHAPE_TYPE.GROUP: GroupShapeExtractor,
}

DEFAULT_EXTRACTOR = BaseShapeExtractor


def shape_extractor_factory(
        shape: Shape, measurement_unit: str = "pt"
) -> ShapeExtractor:
    """Factory function to create a shape extractor based on the shape type or text frame."""
    # First, check if the shape has a text frame
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        logger.info(
            f"Shape ID {shape.shape_id} has a text frame. Using BaseAutoShapeExtractor."
        )
        return BaseAutoShapeExtractor(shape, measurement_unit)

    try:
        shape_type = shape.shape_type
    except NotImplementedError:
        logger.warning(
            f"Unrecognized shape type for shape ID {shape.shape_id}. Using default extractor."
        )
        # Access and log the shape's XML for debugging
        shape_xml = shape.element.xml
        logger.debug(f"Shape ID {shape.shape_id} XML: {shape_xml}")
        return DEFAULT_EXTRACTOR(shape, measurement_unit)

    extractor_class = SHAPE_EXTRACTOR_MAP.get(shape_type, DEFAULT_EXTRACTOR)

    if extractor_class is DEFAULT_EXTRACTOR:
        logger.warning(
            f"No specific extractor for shape type '{shape_type}' (ID: {shape.shape_id}). Using default extractor."
        )

    return extractor_class(shape, measurement_unit)
